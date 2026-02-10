import json
import os
import re
import shutil
import logging
import asyncio
from collections import Counter

import cv2
import numpy as np
from PIL import ImageFont, Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .utils import extract_background_color, extract_font_color, fill_bbox_with_bg, get_projection_segments
from .ai_services import ai_manager, TextCorrectionRequest, TextCorrectionResponse
from .config import ai_config
from .cache_manager import global_cache_manager

logger = logging.getLogger(__name__)


class Character:
    """A simple data class to hold information about a detected character."""

    def __init__(self, bbox, color, line_index):
        self.bbox = bbox
        self.color = color
        self.line_index = line_index
        self.font_size = 0
        self.bold = False
        self.text = ""

    def __repr__(self):
        return (f"Character(text='{self.text}', bbox={self.bbox}, color={self.color}, line={self.line_index}, "
                f"size={self.font_size}, bold={self.bold})")


class PageContext:
    def __init__(self, page_image, coords, slide, page_number=0):
        self.slide = slide
        self.original_image = page_image.copy()
        self.background_image = page_image.copy()
        self.coords = coords
        self.elements = []
        self.raw_chars = []
        self.corrected_chars = []
        self.text_corrections = []  # Store AI correction results
        self.original_texts = {}    # Backup of original texts
        self.page_number = page_number
        
        # Batch text correction data
        self.pending_texts = []     # Texts waiting for correction
        self.text_contexts = []     # Contexts for each text
        self.text_elements = []     # Element references for texts
        # Cache is now managed globally - keeping this for compatibility
        self.corrected_texts_cache = {}  # Legacy cache for backward compatibility

    def clear_corrected_texts_cache(self):
        """Clear the corrected texts cache for this page (now uses global cache)"""
        return global_cache_manager.clear_all()
    
    def get_cache_stats(self):
        """Get statistics about the corrected texts cache (now uses global cache)"""
        return global_cache_manager.get_cache_stats()

    def add_element_bbox_for_cleanup(self, bbox):
        """Register a bounding box to be inpainted on the background image."""
        if bbox:
            px_box = [int(v * (self.coords['img_w'] / self.coords['json_w'] if i % 2 == 0 else self.coords['img_h'] / self.coords['json_h'])) for i, v in enumerate(bbox)]
            fill_bbox_with_bg(self.background_image, px_box)

    def add_processed_element(self, elem_type, data):
        """Store a fully processed element ready for rendering."""
        self.elements.append({'type': elem_type, 'data': data})

    def add_characters(self, raw_chars, corrected_chars):
        if raw_chars: self.raw_chars.extend(raw_chars)
        if corrected_chars: self.corrected_chars.extend(corrected_chars)
    
    def add_text_for_batch_correction(self, text_content: str, context: str = None, element_ref=None):
        """Add text to batch correction queue"""
        self.pending_texts.append(text_content)
        self.text_contexts.append(context or f"Page {self.page_number} text element")
        self.text_elements.append(element_ref)
    
    async def apply_ai_batch_text_correction(self) -> bool:
        """Apply AI text correction to all pending texts in batch with improved management"""
        if not ai_config.is_text_correction_enabled():
            return False
        
        if not self.pending_texts:
            return True
        
        active_provider = ai_config.get_active_provider()
        if not active_provider or not ai_config.is_provider_enabled(active_provider):
            logger.info("No active AI provider for text correction")
            return False
        
        try:
            # Ensure AI manager active service is synchronized with config
            if ai_manager.get_active_service_name() != active_provider:
                if active_provider in ai_manager.services:
                    # Check if the preferred provider is actually authenticated
                    preferred_service = ai_manager.services[active_provider]
                    if preferred_service.is_authenticated:
                        ai_manager.set_active_service_sync(active_provider)
                        logger.info(f"Switched AI manager to active provider: {active_provider}")
                    else:
                        # Preferred provider failed authentication, find any working service
                        logger.warning(f"User's preferred provider {active_provider} not authenticated, looking for alternatives")
                        fallback_provider = None
                        for provider in ai_manager.services:
                            if ai_manager.services[provider].is_authenticated:
                                fallback_provider = provider
                                break
                        
                        if fallback_provider:
                            ai_manager.set_active_service_sync(fallback_provider)
                            logger.warning(f"Using fallback provider {fallback_provider} instead of {active_provider}")
                        else:
                            logger.error("No authenticated AI services available")
                            return False
                else:
                    logger.error(f"AI service {active_provider} not available")
                    return False
            
            # Get the actual active service (might be fallback)
            actual_active = ai_manager.get_active_service_name()
            if not actual_active:
                logger.error("No active AI service available")
                return False
                
            service = ai_manager.services[actual_active]
            
            # Split large batches to prevent truncation and improve reliability
            batches = self._split_texts_into_optimal_batches(
                self.pending_texts, 
                self.text_contexts, 
                max_chars_per_batch=6000  # Conservative limit
            )
            
            all_corrected_texts = []
            
            for batch_idx, (text_batch, context_batch) in enumerate(batches):
                print(f"Processing batch {batch_idx + 1}/{len(batches)} with {len(text_batch)} texts...")
                
                try:
                    # Process batch with size-aware correction
                    corrected_batch = await service.correct_texts_by_page(
                        texts=text_batch,
                        page_number=self.page_number,
                        context=f"Page {self.page_number} content (batch {batch_idx + 1})",
                        correction_type=ai_config.get_correction_type()
                    )
                    
                    # Validate batch response
                    if len(corrected_batch) != len(text_batch):
                        print(f"Warning: Batch {batch_idx + 1} returned {len(corrected_batch)} results, expected {len(text_batch)}")
                        # Pad with original texts if needed
                        while len(corrected_batch) < len(text_batch):
                            missing_index = len(corrected_batch)
                            corrected_batch.append(text_batch[missing_index])
                            print(f"Added missing text {missing_index + 1}")
                    
                    all_corrected_texts.extend(corrected_batch[:len(text_batch)])
                    print(f"Batch {batch_idx + 1} completed successfully")
                    
                except Exception as e:
                    logger.error(f"Error processing batch {batch_idx + 1}: {e}")
                    # Fallback to original texts for this batch
                    all_corrected_texts.extend(text_batch)
                    print(f"Batch {batch_idx + 1} failed, using original texts")
            
            # Cache corrected texts using global cache manager
            active_provider = ai_config.get_active_provider()
            cached_count = 0
            
            for original, corrected in zip(self.pending_texts, all_corrected_texts):
                if corrected and corrected != original:
                    global_cache_manager.cache_correction(original, corrected, active_provider)
                    cached_count += 1
                
                # Store backup if enabled
                if ai_config.should_backup_original():
                    self.original_texts[original] = original
            
            logger.info(f"Batch corrected {len(self.pending_texts)} texts for page {self.page_number} ({cached_count} new corrections cached)")
            
            # Clear pending texts
            self.pending_texts.clear()
            self.text_contexts.clear()
            self.text_elements.clear()
            
            return True
            
        except Exception as e:
            logger.error(f"AI batch text correction failed: {e}")
            return False
    
    async def _process_text_batch_safely(self, text_batch, context_batch, service, batch_idx):
        """Process a single text batch with comprehensive error handling and validation"""
        try:
            # Process batch with size-aware correction
            corrected_batch = await service.correct_texts_by_page(
                texts=text_batch,
                page_number=self.page_number,
                context=f"Page {self.page_number} content (batch {batch_idx + 1})",
                correction_type=ai_config.get_correction_type()
            )
            
            # Comprehensive validation of batch response
            if not corrected_batch:
                logger.warning(f"Batch {batch_idx + 1}: Empty response, using original texts")
                return text_batch
                
            if not isinstance(corrected_batch, list):
                logger.warning(f"Batch {batch_idx + 1}: Response is not a list, using original texts")
                return text_batch
                
            # Handle count mismatch
            if len(corrected_batch) != len(text_batch):
                print(f"Warning: Batch {batch_idx + 1} returned {len(corrected_batch)} results, expected {len(text_batch)}")
                
                # If we got more texts than expected, truncate
                if len(corrected_batch) > len(text_batch):
                    corrected_batch = corrected_batch[:len(text_batch)]
                    print(f"Truncated batch {batch_idx + 1} to {len(text_batch)} texts")
                
                # If we got fewer texts than expected, pad with originals
                while len(corrected_batch) < len(text_batch):
                    missing_index = len(corrected_batch)
                    corrected_batch.append(text_batch[missing_index])
                    print(f"Added missing text {missing_index + 1} from original")
            
            # Validate each corrected text
            validated_batch = []
            for i, (original, corrected) in enumerate(zip(text_batch, corrected_batch)):
                if corrected and isinstance(corrected, str) and corrected.strip():
                    validated_batch.append(corrected)
                else:
                    print(f"Invalid correction for text {i + 1}, using original")
                    validated_batch.append(original)
            
            return validated_batch
            
        except Exception as e:
            logger.error(f"Error processing batch {batch_idx + 1}: {e}")
            # Fallback to original texts for this batch
            print(f"Batch {batch_idx + 1} failed, using original texts")
            return text_batch
    
    def _split_texts_into_optimal_batches(self, texts, contexts, max_chars_per_batch=6000):
        """Split texts into batches that won't exceed AI token limits"""
        batches = []
        current_batch_texts = []
        current_batch_contexts = []
        current_char_count = 0
        
        for text, context in zip(texts, contexts):
            text_chars = len(text) + (len(context) if context else 0)
            
            # If adding this text would exceed the limit and we have texts in current batch
            if current_char_count + text_chars > max_chars_per_batch and current_batch_texts:
                # Start new batch
                batches.append((current_batch_texts, current_batch_contexts))
                current_batch_texts = [text]
                current_batch_contexts = [context]
                current_char_count = text_chars
            else:
                current_batch_texts.append(text)
                current_batch_contexts.append(context)
                current_char_count += text_chars
        
        # Add remaining texts as final batch
        if current_batch_texts:
            batches.append((current_batch_texts, current_batch_contexts))
        
        print(f"Split {len(texts)} texts into {len(batches)} batch(es) (max {max_chars_per_batch} chars each)")
        return batches
    
    async def apply_ai_text_correction(self, text_content: str, context: str = None) -> str:
        """Apply AI text correction to text content (individual fallback)"""
        if not ai_config.is_text_correction_enabled():
            return text_content
        
        # Check global cache first
        active_provider = ai_config.get_active_provider()
        cached_correction = global_cache_manager.get_corrected_text(text_content, active_provider)
        if cached_correction:
            return cached_correction
        
        active_provider = ai_config.get_active_provider()
        if not active_provider or not ai_config.is_provider_enabled(active_provider):
            logger.info("No active AI provider for text correction")
            return text_content
        
        try:
            # Ensure AI manager active service is synchronized with config
            if ai_manager.get_active_service_name() != active_provider:
                if active_provider in ai_manager.services:
                    # Check if the preferred provider is actually authenticated
                    preferred_service = ai_manager.services[active_provider]
                    if preferred_service.is_authenticated:
                        ai_manager.set_active_service_sync(active_provider)
                        logger.info(f"Switched AI manager to active provider for individual correction: {active_provider}")
                    else:
                        # Preferred provider failed authentication, find any working service
                        logger.warning(f"User's preferred provider {active_provider} not authenticated for individual correction, looking for alternatives")
                        fallback_provider = None
                        for provider in ai_manager.services:
                            if ai_manager.services[provider].is_authenticated:
                                fallback_provider = provider
                                break
                        
                        if fallback_provider:
                            ai_manager.set_active_service(fallback_provider)
                            logger.warning(f"Using fallback provider {fallback_provider} for individual correction instead of {active_provider}")
                        else:
                            logger.error("No authenticated AI services available for individual correction")
                            return text_content
                else:
                    logger.error(f"AI service {active_provider} not available for individual correction")
                    return text_content
            
            # Create correction request
            request = TextCorrectionRequest(
                text=text_content,
                context=context,
                language="auto",
                correction_type=ai_config.get_correction_type()
            )
            
            # Apply correction using AI service
            response = await ai_manager.correct_text(request)
            
            # Cache the correction for future use
            global_cache_manager.cache_correction(
                response.original_text,
                response.corrected_text,
                response.provider,
                response.confidence
            )
            
            # Store correction results
            if ai_config.should_backup_original():
                self.original_texts[text_content] = response.original_text
            
            self.text_corrections.append({
                'original': response.original_text,
                'corrected': response.corrected_text,
                'changes': response.changes,
                'confidence': response.confidence,
                'provider': response.provider,
                'model': response.model
            })
            
            logger.info(f"Text corrected by {response.provider} ({response.model}), confidence: {response.confidence:.2f}")
            
            if ai_config.should_show_changes() and response.changes:
                logger.info(f"Changes made: {len(response.changes)} corrections")
                for change in response.changes[:3]:  # Show first 3 changes
                    logger.debug(f"  {change.get('type', 'unknown')}: '{change.get('original', '')}' -> '{change.get('corrected', '')}'")
            
            return response.corrected_text
            
        except Exception as e:
            logger.error(f"AI text correction failed: {e}")
            return text_content  # Return original text if correction fails
    
    def get_correction_summary(self) -> dict:
        """Get summary of all text corrections applied to this page"""
        if not self.text_corrections:
            return {"total_corrections": 0, "providers_used": [], "average_confidence": 0.0}
        
        providers_used = list(set(corr['provider'] for corr in self.text_corrections))
        total_changes = sum(len(corr['changes']) for corr in self.text_corrections)
        avg_confidence = sum(corr['confidence'] for corr in self.text_corrections) / len(self.text_corrections)
        
        return {
            "total_corrections": len(self.text_corrections),
            "total_changes": total_changes,
            "providers_used": providers_used,
            "average_confidence": avg_confidence
        }

    def generate_debug_images(self, page_index, generator_instance):
        """Generate and save debug images for the page."""
        generator_instance._draw_debug_boxes_for_page(self.original_image, self.raw_chars, self.coords, f"tmp/page_{page_index}_raw.png")
        generator_instance._draw_debug_boxes_for_page(self.original_image, self.corrected_chars, self.coords, f"tmp/page_{page_index}_corrected.png")

    def render_to_slide(self, generator_instance):
        """Render all processed elements onto the PowerPoint slide."""
        # 1. Render the cleaned background
        bg_path = f"temp_bg_{id(self.slide)}.png"
        cv2.imwrite(bg_path, cv2.cvtColor(self.background_image, cv2.COLOR_RGB2BGR))
        w_pts, h_pts = generator_instance.prs.slide_width, generator_instance.prs.slide_height
        self.slide.shapes.add_picture(bg_path, Pt(0), Pt(0), w_pts, h_pts)
        os.remove(bg_path)

        # 2. Render all image elements first
        for elem in self.elements:
            if elem['type'] == 'image':
                generator_instance._add_picture_from_bbox(self.slide, elem['data']['bbox'], self.original_image, self.coords, elem['data']['text_elements'])

        # 3. Render all text elements on top
        for elem in self.elements:
            if elem['type'] == 'text':
                # Pass existing elements for collision detection
                generator_instance._render_text_from_data(self.slide, elem['data'], self.elements)


class PPTGenerator:
    def __init__(self, output_path, remove_watermark=True, enable_ai_correction=None):
        self.prs = Presentation()
        self.output_path = output_path
        self.remove_watermark = remove_watermark
        self.debug_images = False # Will be set in process_page
        self.enable_ai_correction = enable_ai_correction if enable_ai_correction is not None else ai_config.is_text_correction_enabled()
        self.ai_initialized = False
        
        # Remove default slide
        for i in range(len(self.prs.slides) - 1, -1, -1):
            rId = self.prs.slides._sldIdLst[i].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[i]
    
    async def initialize_ai_services(self):
        """Initialize AI services based on configuration"""
        if not self.enable_ai_correction:
            logger.info("AI text correction is disabled")
            return False
        
        try:
            # Get enabled providers
            enabled_providers = ai_config.get_enabled_providers()
            if not enabled_providers:
                logger.warning("No AI providers are enabled for text correction")
                return False
            
            # Initialize services for enabled providers
            services_initialized = 0
            for provider in enabled_providers:
                try:
                    api_key = ai_config.get_api_key(provider)
                    model = ai_config.get_model(provider)
                    
                    if not api_key:
                        logger.warning(f"No API key configured for {provider}")
                        continue
                    
                    # Create and register service
                    service = ai_manager.create_service(provider, api_key, model)
                    ai_manager.register_service(provider, service)
                    
                    # Test authentication
                    if await ai_manager.authenticate_service(provider):
                        logger.info(f"Successfully initialized {provider} with model {model}")
                        services_initialized += 1
                    else:
                        logger.error(f"Failed to authenticate with {provider}")
                        
                except Exception as e:
                    logger.error(f"Error initializing {provider}: {e}")
            
            if services_initialized > 0:
                # Set active provider - prioritize user's choice
                active_provider = ai_config.get_active_provider()
                
                # Check if user's chosen provider is available and authenticated
                if active_provider and active_provider in ai_manager.services:
                    # User's choice is available, use it
                    ai_manager.set_active_service(active_provider)
                    logger.info(f"Active AI provider set to user's choice: {active_provider}")
                elif active_provider and active_provider in enabled_providers:
                    # User's choice is enabled but might have failed authentication
                    logger.warning(f"User's chosen provider {active_provider} failed authentication, trying fallback")
                    # Don't override user's choice in config, let them know what happened
                    for provider in enabled_providers:
                        if provider in ai_manager.services:
                            ai_manager.set_active_service(provider)
                            logger.warning(f"Using fallback provider {provider} instead of {active_provider}")
                            break
                else:
                    # No active provider set or invalid choice - use first available
                    for provider in enabled_providers:
                        if provider in ai_manager.services:
                            active_provider = provider
                            ai_config.set_active_provider(provider)
                            ai_manager.set_active_service(provider)
                            logger.info(f"No valid active provider, set to first available: {active_provider}")
                            break
                
                self.ai_initialized = True
                return True
            
            logger.warning("No AI services could be initialized")
            return False
            
        except Exception as e:
            logger.error(f"Error initializing AI services: {e}")
            return False
    
    def clear_all_text_caches(self):
        """Clear all corrected text caches using global cache manager"""
        return global_cache_manager.clear_all()
    
    @staticmethod
    def get_global_cache_stats():
        """Get global statistics about cached texts using global cache manager"""
        return global_cache_manager.get_cache_stats()
    
    @staticmethod
    def clear_cache_by_provider(provider: str):
        """Clear cached corrections for a specific provider"""
        return global_cache_manager.clear_by_provider(provider)
    
    @staticmethod
    def get_recent_corrections(limit: int = 10):
        """Get recent corrections for display"""
        return global_cache_manager.get_recent_corrections(limit)

    def cap_size(self, w_pts, h_pts):
        MAX_PTS = 56 * 72
        if w_pts > MAX_PTS or h_pts > MAX_PTS:
            scale = MAX_PTS / max(w_pts, h_pts)
            w_pts, h_pts = w_pts * scale, h_pts * scale
        return w_pts, h_pts

    def set_slide_size(self, width_px, height_px, dpi=72):
        w_pts, h_pts = self.cap_size(width_px * 72 / dpi, height_px * 72 / dpi)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)

    def add_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _get_bbox_intersection(self, bbox1, bbox2):
        x1, y1 = max(bbox1[0], bbox2[0]), max(bbox1[1], bbox2[1])
        x2, y2 = min(bbox1[2], bbox2[2]), min(bbox1[3], bbox2[3])
        return [x1, y1, x2, y2] if x1 < x2 and y1 < y2 else None

    def _create_textbox(self, slide, bbox, coords):
        x1, y1, x2, y2 = bbox
        return slide.shapes.add_textbox(
            Pt(x1 * coords['scale_x']), Pt(y1 * coords['scale_y']),
            Pt((x2 - x1) * coords['scale_x']), Pt((y2 - y1) * coords['scale_y'])
        )
    
    def _get_bbox_overlap_area(self, bbox1, bbox2):
        """Calculate overlapping area between two bounding boxes"""
        x1_1, y1_1, x2_1, y2_1 = bbox1
        x1_2, y1_2, x2_2, y2_2 = bbox2
        
        # Calculate intersection rectangle
        x_left = max(x1_1, x1_2)
        y_top = max(y1_1, y1_2)
        x_right = min(x2_1, x2_2)
        y_bottom = min(y2_1, y2_2)
        
        # Check if there's an intersection
        if x_right > x_left and y_bottom > y_top:
            return (x_right - x_left) * (y_bottom - y_top)
        return 0
    
    def _find_max_safe_width(self, bbox, existing_elements, margin=5):
        """Find maximum width that won't cause overlaps with existing elements"""
        x1, y1, x2, y2 = bbox
        current_width = x2 - x1
        
        # Find nearest element to the right within the same vertical space
        min_distance = float('inf')
        
        for element in existing_elements:
            elem_data = element.get('data', {})
            elem_bbox = elem_data.get('bbox')
            
            if not elem_bbox:
                continue
            
            elem_x1, elem_y1, elem_x2, elem_y2 = elem_bbox
            
            # Check if element is to the right and has vertical overlap
            if elem_x1 > x2:  # Element is to the right
                # Check for vertical overlap
                vertical_overlap = min(y2, elem_y2) - max(y1, elem_y1)
                if vertical_overlap > 0:  # There's vertical overlap
                    distance = elem_x1 - x2
                    min_distance = min(min_distance, distance)
        
        if min_distance == float('inf'):
            # No elements to the right, can expand more freely
            # But still be conservative - max 50% expansion
            max_safe_width = current_width * 1.5
        else:
            # Leave margin before the next element
            max_safe_width = current_width + (min_distance - margin)
        
        return max_safe_width
    
    def _calculate_optimal_textbox_size(self, bbox, text_content, existing_elements, is_single_line=False):
        """Calculate optimal textbox size without causing overlaps"""
        x1, y1, x2, y2 = bbox
        current_width = x2 - x1
        
        if not is_single_line:
            return bbox  # Don't modify multi-line text boxes
        
        # Calculate content-based width estimate
        if text_content:
            # Rough estimate: average character is about 0.6em wide
            char_count = len(text_content.replace(' ', ''))  # Don't count spaces
            # Estimate width needed (this is very rough)
            estimated_width_factor = max(1.0, char_count * 0.02)  # Rough scaling
            content_based_width = current_width * estimated_width_factor
        else:
            content_based_width = current_width
        
        # Find maximum safe width based on surrounding elements
        max_safe_width = self._find_max_safe_width(bbox, existing_elements)
        
        # Use the smaller of content-based width and safe width, but at least current width
        optimal_width = max(current_width, min(content_based_width, max_safe_width))
        
        # Cap expansion at 50% to prevent excessive growth
        optimal_width = min(optimal_width, current_width * 1.5)
        
        # Only expand if there's meaningful benefit (at least 10% increase)
        if optimal_width > current_width * 1.1:
            new_x2 = x1 + optimal_width
            return [x1, y1, new_x2, y2]
        
        return bbox
    
    def _detect_bullet_points_intelligently(self, text_content, element_type=None, context=None):
        """Enhanced bullet point detection with Hungarian language support and content analysis"""
        if not text_content or not text_content.strip():
            return text_content
        
        # Clean existing bullets first
        cleaned_text = self._clean_existing_bullets(text_content)
        
        lines = cleaned_text.split('\n')
        processed_lines = []
        
        # Analyze overall structure for list detection
        list_indicators = self._analyze_list_structure(lines)
        
        for i, line in enumerate(lines):
            original_line = line
            cleaned_line = line.strip()
            
            if not cleaned_line:
                processed_lines.append(original_line)
                continue
            
            # Multi-stage bullet detection
            should_add_bullet = self._should_add_bullet(
                cleaned_line, element_type, context, list_indicators, i, lines
            )
            
            if should_add_bullet:
                # Preserve original indentation
                indentation = line[:len(line) - len(line.lstrip())]
                processed_lines.append(f"{indentation}• {cleaned_line}")
                print(f"   Intelligent bullet detection: Added bullet to line {i + 1}")
            else:
                processed_lines.append(original_line)
        
        return '\n'.join(processed_lines)
    
    def _clean_existing_bullets(self, text_content):
        """Normalize existing bullet characters to a consistent format"""
        # Replace various bullet characters with standard bullet
        bullet_chars = ['·', '-', '*', '+', '▪', '▫', '◦', '‣', '⁃', '→', '=>']
        
        lines = text_content.split('\n')
        cleaned_lines = []
        
        for line in lines:
            cleaned_line = line
            stripped = line.strip()
            
            # Replace bullet characters at the start of lines
            for bullet_char in bullet_chars:
                if stripped.startswith(bullet_char + ' '):
                    # Preserve indentation
                    indent = line[:len(line) - len(line.lstrip())]
                    content = stripped[2:].strip()  # Remove bullet and space
                    cleaned_line = f"{indent}• {content}"
                    break
            
            cleaned_lines.append(cleaned_line)
        
        return '\n'.join(cleaned_lines)
    
    def _analyze_list_structure(self, lines):
        """Analyze text structure to identify list patterns"""
        indicators = {
            'has_numbered_sequence': False,
            'has_hungarian_ordinals': False,
            'avg_line_length': 0,
            'short_lines_ratio': 0
        }
        
        if not lines:
            return indicators
            
        # Check for numbered sequences
        numbered_lines = 0
        for line in lines:
            stripped = line.strip()
            if re.match(r'^\d+\.?\s+', stripped):
                numbered_lines += 1
        indicators['has_numbered_sequence'] = numbered_lines > 1
        
        # Calculate average length and short lines ratio
        non_empty_lines = [line.strip() for line in lines if line.strip()]
        if non_empty_lines:
            indicators['avg_line_length'] = sum(len(line) for line in non_empty_lines) / len(non_empty_lines)
            short_lines = sum(1 for line in non_empty_lines if len(line) < 60)
            indicators['short_lines_ratio'] = short_lines / len(non_empty_lines)
        
        # Check for Hungarian ordinals
        hungarian_ordinals = ['első', 'második', 'harmadik', 'negyedik', 'ötödik']
        for line in lines:
            if any(ordinal in line.lower() for ordinal in hungarian_ordinals):
                indicators['has_hungarian_ordinals'] = True
                break
        
        return indicators
    
    def _should_add_bullet(self, line, element_type, context, indicators, line_index, all_lines):
        """Determine if a line should have a bullet point added"""
        # Already has a bullet
        if line.startswith('• '):
            return False
            
        # Skip headers (too long, ends with colon, or contains specific patterns)
        if (len(line) > 100 or 
            line.endswith(':') or 
            any(pattern in line.lower() for pattern in ['cím:', 'fejezet:', 'rész:', 'típus:'])):
            return False
        
        # Element type context
        if element_type in ['list', 'para_blocks']:
            return True
            
        # Numbered sequence detection
        if indicators['has_numbered_sequence'] and re.match(r'^\d+\.?\s+', line):
            return True
            
        # Hungarian ordinal detection
        if indicators['has_hungarian_ordinals']:
            hungarian_starters = ['első', 'második', 'harmadik', 'negyedik', 'ötödik', 
                                'hatodik', 'hetedik', 'nyolcadik', 'kilencedik', 'tizedik']
            if any(line.lower().startswith(starter) for starter in hungarian_starters):
                return True
        
        # Hungarian-specific bullet indicators
        hungarian_patterns = [
            line.startswith(('Például', 'például')),     # "For example"
            line.startswith(('Amely', 'amely')),         # "Which"
            line.startswith(('Amikor', 'amikor')),       # "When"
            line.startswith(('Ahogy', 'ahogy')),         # "As"
            line.startswith(('Mivel', 'mivel')),         # "Since"
            line.startswith(('Nem', 'nem')),             # "Not"
            any(word in line.lower() for word in ['továbbá', 'emellett', 'valamint']),  # Connectors
        ]
        
        # General bullet patterns
        general_patterns = [
            line[0].islower() and len(line) > 15,        # Starts lowercase, substantial
            line.count(',') >= 2,                         # Multiple clauses
            line.endswith(('.', ':', ';', '!')),         # Typical endings
            20 < len(line) < 80,                         # Reasonable bullet length
        ]
        
        hungarian_score = sum(hungarian_patterns)
        general_score = sum(general_patterns)
        
        # Decision logic: be more conservative to avoid false positives
        return hungarian_score >= 1 or general_score >= 3
    
    def _clean_existing_bullets(self, text_content):
        """Clean and normalize existing bullet characters"""
        if not text_content:
            return text_content
        
        # Replace various bullet characters with consistent bullet
        bullet_chars = ['·', '*', '-', '○', '■', '▪', '▫', '◦']
        
        lines = text_content.split('\n')
        cleaned_lines = []
        
        for line in lines:
            stripped = line.lstrip()
            leading_space = line[:len(line) - len(stripped)]
            
            # Replace bullet characters at the start
            for bullet_char in bullet_chars:
                if stripped.startswith(bullet_char):
                    # Remove old bullet and clean spacing
                    remaining = stripped[1:].lstrip()
                    if remaining:  # Only if there's content after the bullet
                        stripped = "• " + remaining
                    break
            
            cleaned_lines.append(leading_space + stripped)
        
        return '\n'.join(cleaned_lines)
    
    def _enhanced_character_alignment(self, raw_chars, full_text, coords):
        """Enhanced character alignment with better fallback handling"""
        
        # Original alignment attempt
        corrected_chars = self._analyze_and_correct_bboxes(raw_chars, full_text, coords)
        non_space_chars = [c for c in full_text if c not in " \n\t"]
        
        # Primary alignment check
        primary_alignment = len(corrected_chars) == len(non_space_chars)
        
        if primary_alignment:
            print(f"Character alignment successful: Primary alignment (1:1 mapping)")
            return corrected_chars, True
        
        print(f"Primary alignment failed: {len(corrected_chars)} chars detected vs {len(non_space_chars)} expected")
        
        # Secondary alignment: try with character merging
        if len(corrected_chars) > len(non_space_chars):
            print(f"Attempting character merging ({len(corrected_chars)} -> {len(non_space_chars)})...")
            merged_chars = self._merge_close_characters(corrected_chars, coords)
            if len(merged_chars) == len(non_space_chars):
                print(f"Character alignment successful: Merging fixed alignment")
                return merged_chars, True
        
        # Tertiary alignment: try with character splitting
        if len(corrected_chars) < len(non_space_chars):
            print(f"Attempting character splitting ({len(corrected_chars)} -> {len(non_space_chars)})...")
            split_chars = self._split_wide_characters(corrected_chars, non_space_chars, coords)
            if len(split_chars) == len(non_space_chars):
                print(f"Character alignment successful: Splitting fixed alignment")
                return split_chars, True
        
        # Quaternary alignment: try intelligent character mapping
        print(f"Attempting intelligent fallback mapping...")
        fallback_chars = self._create_fallback_character_mapping(raw_chars, full_text, coords)
        
        if len(fallback_chars) == len(non_space_chars):
            print(f"Character alignment successful: Intelligent fallback mapping")
            return fallback_chars, True
        
        # If all else fails, return best effort
        print(f"Character alignment failed: Using best available mapping ({len(corrected_chars)} chars)")
        return corrected_chars, False

    def _merge_close_characters(self, chars, coords):
        """Merge characters that are very close together (likely over-segmented)"""
        if not chars:
            return chars
            
        merged = []
        scale_x = coords.get('scale_x', 1)
        i = 0
        
        while i < len(chars):
            current_char = chars[i]
            
            # Look ahead for very close characters
            if i + 1 < len(chars):
                next_char = chars[i + 1]
                
                # Calculate distance between characters
                current_right = current_char.bbox[2]
                next_left = next_char.bbox[0]
                distance = next_left - current_right
                
                # Merge threshold: very close characters (less than 3 pixels scaled)
                merge_threshold = 3.0 / scale_x if scale_x > 0 else 3.0
                
                if distance <= merge_threshold and current_char.line_index == next_char.line_index:
                    # Create merged character
                    merged_bbox = [
                        current_char.bbox[0],
                        min(current_char.bbox[1], next_char.bbox[1]),
                        next_char.bbox[2],
                        max(current_char.bbox[3], next_char.bbox[3])
                    ]
                    
                    # Use the color of the larger character
                    current_area = (current_char.bbox[2] - current_char.bbox[0]) * (current_char.bbox[3] - current_char.bbox[1])
                    next_area = (next_char.bbox[2] - next_char.bbox[0]) * (next_char.bbox[3] - next_char.bbox[1])
                    dominant_color = current_char.color if current_area >= next_area else next_char.color
                    
                    merged_char = Character(merged_bbox, dominant_color, current_char.line_index)
                    merged.append(merged_char)
                    i += 2  # Skip next character as it's been merged
                    continue
            
            # No merge possible, add character as-is
            merged.append(current_char)
            i += 1
        
        return merged
    
    def _split_wide_characters(self, chars, expected_chars, coords):
        """Split wide characters that likely represent multiple characters"""
        if not chars:
            return chars
            
        chars_needed = len(expected_chars)
        chars_available = len(chars)
        
        if chars_available >= chars_needed:
            return chars  # No splitting needed
        
        split_chars = []
        chars_to_generate = chars_needed - chars_available
        scale_x = coords.get('scale_x', 1)
        
        # Sort characters by width to identify candidates for splitting
        char_widths = [(i, char, char.bbox[2] - char.bbox[0]) for i, char in enumerate(chars)]
        char_widths.sort(key=lambda x: x[2], reverse=True)  # Widest first
        
        # Split the widest characters until we have enough
        split_candidates = char_widths[:min(chars_to_generate, len(char_widths))]
        
        for i, char in enumerate(chars):
            # Check if this character should be split
            should_split = any(candidate[1] == char for candidate in split_candidates)
            
            if should_split and chars_to_generate > 0:
                # Split this character into 2 parts
                char_width = char.bbox[2] - char.bbox[0]
                split_point = char.bbox[0] + char_width / 2
                
                # Create two characters from one
                left_char = Character(
                    [char.bbox[0], char.bbox[1], split_point, char.bbox[3]],
                    char.color,
                    char.line_index
                )
                right_char = Character(
                    [split_point, char.bbox[1], char.bbox[2], char.bbox[3]],
                    char.color,
                    char.line_index
                )
                
                split_chars.extend([left_char, right_char])
                chars_to_generate -= 1
            else:
                split_chars.append(char)
        
        return split_chars
    
    def _create_fallback_character_mapping(self, raw_chars, full_text, coords):
        """Create fallback character mapping when alignment fails"""
        non_space_chars = [c for c in full_text if c not in " \n\t"]
        
        if not raw_chars or not non_space_chars:
            return raw_chars
        
        # Create synthetic characters based on text length and available space
        if raw_chars:
            # Use the bounding area of all raw characters
            min_x = min(char.bbox[0] for char in raw_chars)
            min_y = min(char.bbox[1] for char in raw_chars)
            max_x = max(char.bbox[2] for char in raw_chars)
            max_y = max(char.bbox[3] for char in raw_chars)
            
            total_width = max_x - min_x
            char_count = len(non_space_chars)
            
            if char_count > 0:
                char_width = total_width / char_count
                fallback_chars = []
                
                # Get the most common color and line index
                colors = [char.color for char in raw_chars]
                line_indices = [char.line_index for char in raw_chars]
                common_color = max(set(colors), key=colors.count) if colors else (0, 0, 0)
                common_line = max(set(line_indices), key=line_indices.count) if line_indices else 0
                
                # Create evenly distributed characters
                for i in range(char_count):
                    char_left = min_x + i * char_width
                    char_right = min_x + (i + 1) * char_width
                    
                    char_bbox = [char_left, min_y, char_right, max_y]
                    fallback_char = Character(char_bbox, common_color, common_line)
                    fallback_chars.append(fallback_char)
                
                return fallback_chars
        
        return raw_chars
    
    def _detect_language(self, text):
        """Simple language detection based on character patterns"""
        if not text or len(text.strip()) == 0:
            return 'default'
        
        # Hungarian-specific characters
        hungarian_chars = set('áéíóöőúüűÁÉÍÓÖŐÚÜŰ')
        
        char_count = len(text)
        hungarian_char_count = sum(1 for c in text if c in hungarian_chars)
        
        # If more than 1.5% of characters are Hungarian-specific
        if char_count > 0 and hungarian_char_count / char_count > 0.015:
            return 'hungarian'
        
        # Hungarian word patterns
        hungarian_words = [
            'hogy', 'amely', 'amikor', 'mivel', 'tehát', 'valamint', 'továbbá',
            'például', 'között', 'után', 'előtt', 'alatt', 'felett', 'mellett',
            'évek', 'során', 'után', 'előtte', 'utána'
        ]
        
        text_lower = text.lower()
        hungarian_word_count = sum(1 for word in hungarian_words if word in text_lower)
        
        if hungarian_word_count >= 2:
            return 'hungarian'
        
        # German detection (for scientific content)
        german_chars = set('äöüßÄÖÜ')
        german_char_count = sum(1 for c in text if c in german_chars)
        german_words = ['der', 'die', 'das', 'und', 'oder', 'mit', 'für', 'von', 'zu', 'im', 'am', 'auf']
        german_word_count = sum(1 for word in german_words if word in text_lower)
        
        if (german_char_count > 0 and german_word_count >= 2) or german_word_count >= 3:
            return 'german'
        
        return 'default'
    
    def _select_optimal_font_for_language(self, text_content, detected_language=None):
        """Select optimal font based on text content and language"""
        
        # Detect language if not provided
        if not detected_language:
            detected_language = self._detect_language(text_content)
        
        # Hungarian-specific font preferences (excellent Unicode support)
        if detected_language == 'hungarian':
            hungarian_fonts = [
                "Times New Roman",  # Excellent Hungarian support, professional
                "Arial",           # Good fallback, wide compatibility
                "Calibri",         # Modern, clean, good Unicode support
                "Segoe UI",        # Windows standard with excellent Unicode
                "Georgia",         # Good for body text, excellent readability
                "Verdana"          # Clear at small sizes
            ]
            return hungarian_fonts[0]
        
        # German text optimization
        elif detected_language == 'german':
            german_fonts = [
                "Times New Roman",  # Traditional, excellent German support
                "Arial",           # Clean, professional
                "Calibri"          # Modern alternative
            ]
            return german_fonts[0]
        
        # English and other languages
        language_fonts = {
            'english': "Calibri",     # Clean, modern
            'french': "Arial",        # Good accented character support
            'spanish': "Arial",       # Good accented character support
            'italian': "Times New Roman",  # Classical, good readability
            'default': "Arial"        # Safe fallback with good Unicode support
        }
        
        return language_fonts.get(detected_language, language_fonts['default'])

    def _get_line_ranges(self, page_image, bbox, coords):
        x1, y1, x2, y2 = bbox
        px1, py1, px2, py2 = int(x1 * coords['img_w'] / coords['json_w']), int(
            y1 * coords['img_h'] / coords['json_h']), int(x2 * coords['img_w'] / coords['json_w']), int(
            y2 * coords['img_h'] / coords['json_h'])
        h, w = page_image.shape[:2];
        px1, py1, px2, py2 = max(0, px1), max(0, py1), min(w, px2), min(h, py2)
        if px2 <= px1 or py2 <= py1: return []
        roi = page_image[py1:py2, px1:px2]
        bg_color = extract_background_color(page_image, [px1, py1, px2, py2])
        font_color, _, _ = extract_font_color(page_image, [px1, py1, px2, py2], bg_color)
        initial_lines = get_projection_segments(roi, font_color, axis=1)
        line_infos = []
        scale_y = (y2 - y1) / roi.shape[0] if roi.shape[0] > 0 else 0
        for start_y, end_y in initial_lines:
            line_pixel_bbox = [px1, py1 + start_y, px2, py1 + end_y]
            line_bg = extract_background_color(page_image, line_pixel_bbox)
            line_fg, _, _ = extract_font_color(page_image, line_pixel_bbox, line_bg)
            line_infos.append({'range': [y1 + start_y * scale_y, y1 + end_y * scale_y], 'color': line_fg,
                               'pixel_range': (start_y, end_y)})
        if not line_infos: return []
        avg_line_height = np.mean([info['pixel_range'][1] - info['pixel_range'][0] for info in line_infos])
        recovered_lines = []
        sorted_lines = sorted(line_infos, key=lambda x: x['pixel_range'][0])
        all_gaps = [(0, sorted_lines[0]['pixel_range'][0])] + [
            (sorted_lines[i]['pixel_range'][1], sorted_lines[i + 1]['pixel_range'][0]) for i in
            range(len(sorted_lines) - 1)] + [(sorted_lines[-1]['pixel_range'][1], roi.shape[0])]
        for gap_start, gap_end in all_gaps:
            if (gap_end - gap_start) > avg_line_height * 0.8:
                gap_bbox = [px1, py1 + gap_start, px2, py1 + gap_end]
                gap_bg = extract_background_color(page_image, gap_bbox)
                new_font_color, x_prop, y_prop = extract_font_color(page_image, gap_bbox, gap_bg)
                if y_prop > x_prop * 1.2 and np.linalg.norm(np.array(new_font_color) - np.array(font_color)) > 50:
                    gap_roi = roi[gap_start:gap_end, :];
                    gap_pixels = gap_roi.reshape(-1, 3)
                    gap_diff = np.linalg.norm(gap_pixels - new_font_color, axis=1)
                    gap_mask = (gap_diff < 40).reshape(gap_roi.shape[:2])
                    gap_row_counts = np.sum(gap_mask, axis=1)
                    in_gap_line, gap_line_start = False, 0
                    for y, count in enumerate(gap_row_counts):
                        if count > 1 and not in_gap_line:
                            in_gap_line, gap_line_start = True, y
                        elif count < 1 and in_gap_line:
                            in_gap_line = False
                            if y - gap_line_start > 3:
                                abs_start, abs_end = gap_start + gap_line_start, gap_start + y
                                recovered_lines.append({'range': [y1 + abs_start * scale_y, y1 + abs_end * scale_y],
                                                        'color': new_font_color, 'pixel_range': (abs_start, abs_end)})
                    if in_gap_line:
                        abs_start, abs_end = gap_start + gap_line_start, gap_start + len(gap_row_counts)
                        recovered_lines.append(
                            {'range': [y1 + abs_start * scale_y, y1 + abs_end * scale_y], 'color': new_font_color,
                             'pixel_range': (abs_start, abs_end)})
        if recovered_lines:
            line_infos.extend(recovered_lines)
            line_infos.sort(key=lambda x: x['range'][0])

        if len(line_infos) > 1:
            avg_line_height = np.mean([info['pixel_range'][1] - info['pixel_range'][0] for info in line_infos])
            merged_lines = [line_infos[0]]
            for i in range(1, len(line_infos)):
                prev_line = merged_lines[-1]
                curr_line = line_infos[i]
                gap = curr_line['pixel_range'][0] - prev_line['pixel_range'][1]
                if gap >= 0 and gap <= max(avg_line_height * 0.05, 1):
                    prev_height = prev_line['pixel_range'][1] - prev_line['pixel_range'][0]
                    curr_height = curr_line['pixel_range'][1] - curr_line['pixel_range'][0]
                    new_pixel_range = (prev_line['pixel_range'][0], curr_line['pixel_range'][1])
                    new_range = [prev_line['range'][0], curr_line['range'][1]]
                    new_color = curr_line['color'] if curr_height > prev_height else prev_line['color']
                    merged_lines[-1] = {'range': new_range, 'color': new_color, 'pixel_range': new_pixel_range}
                else:
                    merged_lines.append(curr_line)
            line_infos = merged_lines

        for i, info in enumerate(line_infos):
            # Define the top of the search area as the bottom of the previous line.
            info['search_top_y'] = line_infos[i - 1]['range'][1] if i > 0 else bbox[1]

        return line_infos

    def _detect_raw_characters(self, page_image, line_infos, bbox, coords):
        char_objects = []
        for i, info in enumerate(line_infos):
            tight_bbox = [bbox[0], info['range'][0], bbox[2], info['range'][1]]
            search_top_y = info['search_top_y']
            char_objects.extend(
                self._detect_characters_from_line(page_image, tight_bbox, search_top_y, coords, info['color'], i))
        return char_objects

    def _detect_characters_from_line(self, page_image, tight_bbox, search_top_y, coords, line_color, line_index):
        x1, y1, x2, y2 = tight_bbox
        # Convert JSON coordinates to pixel coordinates for the tight box and the search boundary
        px1 = int(x1 * coords['img_w'] / coords['json_w'])
        py1 = int(y1 * coords['img_h'] / coords['json_h'])
        px2 = int(x2 * coords['img_w'] / coords['json_w'])
        py2 = int(y2 * coords['img_h'] / coords['json_h'])
        search_top_py = int(search_top_y * coords['img_h'] / coords['json_h'])

        h, w = page_image.shape[:2]
        px1, py1, px2, py2 = max(0, px1), max(0, py1), min(w, px2), min(h, py2)
        search_top_py = max(0, search_top_py)

        if px2 <= px1 or py2 <= py1:
            return []

        # Define the single, consistent scaling factors for this line based on the tight box.
        scale_x = (x2 - x1) / (px2 - px1) if (px2 - px1) > 0 else 0
        scale_y = (y2 - y1) / (py2 - py1) if (py2 - py1) > 0 else 0

        # Define the region of interest strictly for the tight line box for primary character segmentation.
        tight_roi = page_image[py1:py2, px1:px2]

        # Segment main characters within the TIGHT ROI.
        all_chars = self._segment_characters_in_roi(
            tight_roi, tight_bbox, line_color, line_index, scale_x
        )

        if not all_chars:
            return []
        sorted_chars = sorted(all_chars, key=lambda c: c.bbox[0])

        # Find gaps between characters to search for text of a different color.
        gaps, last_x2 = [], tight_bbox[0]
        for char in sorted_chars:
            if char.bbox[0] > last_x2: gaps.append((last_x2, char.bbox[0]))
            last_x2 = char.bbox[2]
        if tight_bbox[2] > last_x2: gaps.append((last_x2, tight_bbox[2]))

        recovered_chars = []
        scale_x_inv = (px2 - px1) / (x2 - x1) if (x2 - x1) > 0 else 0
        for gap_x1, gap_x2 in gaps:
            gap_px1 = px1 + int((gap_x1 - x1) * scale_x_inv)
            gap_px2 = px1 + int((gap_x2 - x1) * scale_x_inv)

            # Shrink the search box by 5 pixels on each side to avoid edge artifacts from the primary font.
            gap_px1 += 5
            gap_px2 -= 5

            if gap_px2 - gap_px1 < 30: continue
            gap_roi = page_image[search_top_py:py2, gap_px1:gap_px2]
            if gap_roi.size == 0: continue

            gap_bg = extract_background_color(page_image, [gap_px1, search_top_py, gap_px2, py2])
            new_font_color, x_prop, y_prop = extract_font_color(page_image,
                                                                [gap_px1, search_top_py, gap_px2, py2], gap_bg)

            if max(x_prop, y_prop) > 0.15 and np.linalg.norm(np.array(new_font_color) - np.array(line_color)) > 50:
                segments = get_projection_segments(gap_roi, new_font_color, axis=1)

                if segments:
                    # Find the tallest segment, as it's the most likely candidate for the actual line of text.
                    best_segment = max(segments, key=lambda s: s[1] - s[0])
                    segment_height = best_segment[1] - best_segment[0]

                    if segment_height >= 8:
                        local_py1 = best_segment[0]
                        adjusted_roi_py1 = search_top_py + local_py1
                        adjusted_gap_roi = page_image[adjusted_roi_py1:py2, gap_px1:gap_px2]

                        if adjusted_gap_roi.size > 0:
                            new_tight_y1 = search_top_y + (adjusted_roi_py1 - search_top_py) * scale_y
                            new_tight_bbox = [gap_x1, new_tight_y1, gap_x2, y2]
                            recovered_chars.extend(
                                self._segment_characters_in_roi(
                                    adjusted_gap_roi, new_tight_bbox, new_font_color, line_index, scale_x
                                )
                            )
        if recovered_chars:
            all_chars.extend(recovered_chars)
            all_chars.sort(key=lambda c: c.bbox[0])
        return all_chars

    def _segment_characters_in_roi(self, roi, tight_bbox, color, line_index, scale_x):
        # Detect horizontal pixel segments where characters are likely present.
        character_pixel_segments = get_projection_segments(roi, color, axis=0, min_length=2)
        if not character_pixel_segments:
            return []

        char_objects = []
        min_char_pixel_width = 4

        # Unpack the tight bounding box in JSON coordinates.
        json_x1, json_y1, _, json_y2 = tight_bbox

        # Iterate over each detected pixel segment.
        for start_px, end_px in character_pixel_segments:
            if end_px - start_px < min_char_pixel_width:
                continue

            # Convert the relative pixel coordinates of the segment back to absolute JSON coordinates.
            char_json_x1 = json_x1 + start_px * scale_x
            char_json_x2 = json_x1 + end_px * scale_x

            # Create the final bounding box for the character.
            char_bbox = [char_json_x1, json_y1, char_json_x2, json_y2]
            char_objects.append(Character(bbox=char_bbox, color=color, line_index=line_index))

        return char_objects

    def _analyze_and_correct_bboxes(self, char_objects, full_text, coords):
        non_space_chars = [c for c in full_text if c not in " \n"]
        expected_count = len(non_space_chars)
        if not char_objects or len(char_objects) < expected_count:
            return char_objects

        chars = sorted(char_objects, key=lambda c: (c.line_index, c.bbox[0]))
        num_fragments = len(chars)
        num_chars = len(non_space_chars)

        try:
            font = ImageFont.truetype("msyh.ttc", size=30)
            ideal_height = 30
            ideal_char_ratios = []
            for c in non_space_chars:
                if c in "iI,":
                    ideal_char_ratios.append(0.15)
                elif c in "，。、；：？！（）‘’":
                    ideal_char_ratios.append(0.25)
                elif c in "【】“”《》":
                    ideal_char_ratios.append(0.35)
                else:
                    ideal_char_ratios.append(font.getlength(c) / ideal_height)
        except IOError:
            return chars

        memo_cost = {}

        def get_merge_cost(start, end, char_idx):
            if (start, end, char_idx) in memo_cost:
                return memo_cost[(start, end, char_idx)]

            # Rule 2: Don't merge boxes of different colors.
            first_color = chars[start].color
            for i in range(start + 1, end):
                if chars[i].color != first_color:
                    memo_cost[(start, end, char_idx)] = float('inf')
                    return float('inf')

            # Rule 1: Don't merge boxes of different heights (with 20% tolerance).
            heights = [c.bbox[3] - c.bbox[1] for c in chars[start:end]]
            if not heights:
                memo_cost[(start, end, char_idx)] = float('inf')
                return float('inf')
            min_h, max_h = min(heights), max(heights)
            if max_h > min_h * 1.2:
                memo_cost[(start, end, char_idx)] = float('inf')
                return float('inf')

            merged_bbox = [chars[start].bbox[0],
                           min(c.bbox[1] for c in chars[start:end]),
                           chars[end - 1].bbox[2],
                           max(c.bbox[3] for c in chars[start:end])]

            merged_width = merged_bbox[2] - merged_bbox[0]
            merged_height = merged_bbox[3] - merged_bbox[1]
            if merged_height == 0:
                return float('inf')

            merged_ratio = merged_width / merged_height
            ideal_ratio = ideal_char_ratios[char_idx]

            # If ideal ratio is for a full-width char (close to 1) and the detected ratio
            # is slightly narrower (0.9-1.0), treat it as a perfect match (cost 0).
            if ideal_ratio > 0.9 and 0.9 <= merged_ratio <= 1.0:
                cost = 0
            else:
                cost = abs(merged_ratio - ideal_ratio)

            # Adjust cost based on the gap with the preceding character. A larger gap reduces the cost.
            gap_width = 0
            if start > 0:
                gap = chars[start].bbox[0] - chars[start - 1].bbox[2]
                if gap > 0:
                    gap_width = gap

            if merged_height > 0:
                cost -= 0.1 * (gap_width / merged_height)

            memo_cost[(start, end, char_idx)] = cost
            return cost

        dp = [[float('inf')] * (num_chars + 1) for _ in range(num_fragments + 1)]
        path = [[0] * (num_chars + 1) for _ in range(num_fragments + 1)]
        dp[0][0] = 0

        for j in range(1, num_chars + 1):
            for i in range(1, num_fragments + 1):
                for k in range(i):
                    cost = get_merge_cost(k, i, j - 1)
                    if dp[k][j - 1] + cost < dp[i][j]:
                        dp[i][j] = dp[k][j - 1] + cost
                        path[i][j] = k

        if dp[num_fragments][num_chars] == float('inf'):
            return []

        final_chars = []
        curr_frag = num_fragments
        for curr_char in range(num_chars, 0, -1):
            prev_frag = path[curr_frag][curr_char]

            merged_bbox = [chars[prev_frag].bbox[0],
                           min(c.bbox[1] for c in chars[prev_frag:curr_frag]),
                           chars[curr_frag - 1].bbox[2],
                           max(c.bbox[3] for c in chars[prev_frag:curr_frag])]

            new_char = Character(merged_bbox, chars[prev_frag].color, chars[prev_frag].line_index)
            new_char.text = non_space_chars[curr_char - 1]
            final_chars.append(new_char)
            curr_frag = prev_frag

        return final_chars[::-1]

    def _normalize_font_sizes(self, styles):
        if not styles:
            return styles

        i = 0
        while i < len(styles):
            j = i
            while j + 1 < len(styles) and abs(styles[j + 1].font_size - styles[j].font_size) < 3:
                j += 1

            group = styles[i:j + 1]
            if group:
                sizes = [s.font_size for s in group]
                most_common_size = Counter(sizes).most_common(1)[0][0]
                for style in group:
                    style.font_size = most_common_size

            i = j + 1

        return styles

    def _normalize_colors(self, styles, threshold=50):
        if not styles:
            return styles

        i = 0
        while i < len(styles):
            j = i
            while j + 1 < len(styles) and np.linalg.norm(
                    np.array(styles[j + 1].color) - np.array(styles[j].color)) < threshold:
                j += 1

            group = styles[i:j + 1]
            if group:
                colors = [tuple(s.color) for s in group]
                most_common_color = Counter(colors).most_common(1)[0][0]
                for style in group:
                    style.color = most_common_color
            i = j + 1
        return styles

    def _determine_character_styles(self, final_chars, coords, elem_type):
        for char in final_chars:
            height_pts = (char.bbox[3] - char.bbox[1]) * 0.95 * coords['scale_y']
            char.font_size = int(max(height_pts, 6.0))
            char.bold = elem_type == "title"
        return final_chars

    def _draw_debug_boxes_for_page(self, image, all_chars, coords, output_path):
        """Draws bounding boxes for an entire page's characters for debugging."""
        debug_img = image.copy()
        for char in all_chars:
            bbox = char.bbox
            px_box = [
                int(bbox[0] * coords['img_w'] / coords['json_w']),
                int(bbox[1] * coords['img_h'] / coords['json_h']),
                int(bbox[2] * coords['img_w'] / coords['json_w']),
                int(bbox[3] * coords['img_h'] / coords['json_h'])
            ]
            cv2.rectangle(debug_img, (px_box[0], px_box[1]), (px_box[2], px_box[3]), (0, 0, 255), 2)  # Red box
        cv2.imwrite(output_path, cv2.cvtColor(debug_img, cv2.COLOR_RGB2BGR))

    def _process_text(self, context, elem):
        bbox = elem.get("bbox")
        if not bbox: return

        context.add_element_bbox_for_cleanup(bbox)

        all_spans = [s for l in elem.get("lines", []) for s in l.get("spans", [])] if "lines" in elem else elem.get(
            "spans", [])
        if not all_spans:
            text_content = elem.get("text", "")
            if text_content:
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': [{'text': text_content}]})
            return

        # Pre-process spans for bullet points with intelligent detection
        if all_spans:
            first_span_content = all_spans[0].get("content", "")
            if first_span_content.lstrip().startswith('-'):
                cleaned_content = first_span_content.lstrip(' \t\n\r\f\v-•*·')
                all_spans[0]["content"] = "• " + cleaned_content

        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in all_spans])
        if not full_text.strip(): return
        print(f"\n--- Processing Text ---\nContent: '{full_text.strip()[:100]}...'")

        # Apply intelligent bullet point detection
        original_full_text = full_text
        full_text = self._detect_bullet_points_intelligently(full_text, elem.get("type"), context)
        if full_text != original_full_text:
            # Update spans with the new text
            self._update_spans_with_corrected_text(all_spans, original_full_text, full_text)

        try:
            line_infos = self._get_line_ranges(context.original_image, bbox, context.coords)
            print(f"Detected lines: {len(line_infos)}")
            if not line_infos: raise ValueError("No lines detected.")

            raw_chars = self._detect_raw_characters(context.original_image, line_infos, bbox, context.coords)

            # Removed old color-based bullet heuristic - replaced with intelligent detection above

            # Use enhanced character alignment system for better reliability
            corrected_chars, can_align = self._enhanced_character_alignment(raw_chars, full_text, context.coords)
            context.add_characters(raw_chars, corrected_chars)

            print(f"Using char-by-char styling (enhanced alignment): {can_align}")

            if not can_align:
                # Fallback to line-based rendering
                text_runs = self._get_text_runs_by_line(all_spans, line_infos, context.coords, elem.get("type"))
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs})
                return

            final_styles = self._determine_character_styles(corrected_chars, context.coords, elem.get("type"))
            final_styles = self._normalize_font_sizes(final_styles)
            final_styles = self._normalize_colors(final_styles)

            # Determine if single-line based on character analysis BEFORE adding to context
            line_indices = {char.line_index for char in final_styles}
            is_single_line = len(line_indices) <= 1

            text_runs = self._get_text_runs_by_char(full_text, final_styles)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})

        except Exception:
            # Broad exception fallback
            text_runs = self._get_text_runs_from_spans(all_spans, bbox, context.original_image, context.coords, elem.get("type"))
            # Fallback check for single-line
            is_single_line = not any('\n' in run['text'] for run in text_runs)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})
    
    async def _process_text_async(self, context, elem):
        """Async version of text processing with AI text correction"""
        bbox = elem.get("bbox")
        if not bbox: return

        context.add_element_bbox_for_cleanup(bbox)

        all_spans = [s for l in elem.get("lines", []) for s in l.get("spans", [])] if "lines" in elem else elem.get(
            "spans", [])
        if not all_spans:
            text_content = elem.get("text", "")
            if text_content:
                # Use cached corrected text or fallback to individual correction
                if self.enable_ai_correction and self.ai_initialized:
                    active_provider = ai_config.get_active_provider()
                    cached_correction = global_cache_manager.get_corrected_text(text_content, active_provider)
                    if cached_correction:
                        text_content = cached_correction
                        print(f"Using cached correction for simple text: '{text_content[:50]}...'")
                    else:
                        text_content = await context.apply_ai_text_correction(text_content, f"Element type: {elem.get('type', 'text')}")
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': [{'text': text_content}]})
            return

        # Pre-process spans for bullet points with intelligent detection
        if all_spans:
            first_span_content = all_spans[0].get("content", "")
            if first_span_content.lstrip().startswith('-'):
                cleaned_content = first_span_content.lstrip(' \t\n\r\f\v-•*·')
                all_spans[0]["content"] = "• " + cleaned_content

        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in all_spans])
        if not full_text.strip(): return
        
        # Apply intelligent bullet point detection before AI processing
        original_full_text = full_text
        bullet_corrected_text = self._detect_bullet_points_intelligently(full_text, elem.get("type"), context)
        if bullet_corrected_text != full_text:
            full_text = bullet_corrected_text
            # Update spans with bullet-corrected text
            self._update_spans_with_corrected_text(all_spans, original_full_text, full_text)
            original_full_text = full_text  # Update reference for AI correction
        
        # Use cached corrected text from batch processing
        if self.enable_ai_correction and self.ai_initialized:
            # Check cache for corrected text
            active_provider = ai_config.get_active_provider()
            cached_correction = global_cache_manager.get_corrected_text(full_text, active_provider)
            if cached_correction and cached_correction != full_text:
                print(f"Using cached AI correction: '{full_text[:50]}...' -> '{cached_correction[:50]}...'")
                full_text = cached_correction
                
                # Update spans with corrected content if possible
                if len(all_spans) == 1:
                    # Simple case: single span
                    all_spans[0]["content"] = cached_correction
                else:
                    # Complex case: multiple spans - distribute corrected text proportionally
                    self._update_spans_with_corrected_text(all_spans, full_text, cached_correction)
            else:
                # Fallback to individual correction if not in cache
                elem_type = elem.get("type", "text")
                context_info = f"Element type: {elem_type}"
                if elem_type in ["title", "header"]:
                    context_info += " (heading/title text)"
                elif elem_type in ["caption", "footnote"]:
                    context_info += " (caption/footnote text)"
                
                corrected_text = await context.apply_ai_text_correction(full_text, context_info)
                if corrected_text != full_text:
                    print(f"Individual AI Text Correction: '{full_text[:50]}...' -> '{corrected_text[:50]}...'")
                    full_text = corrected_text
                    
                    # Update spans with corrected content
                    if len(all_spans) == 1:
                        all_spans[0]["content"] = corrected_text
                    else:
                        self._update_spans_with_corrected_text(all_spans, original_full_text, corrected_text)
        
        print(f"\n--- Processing Text ---\nContent: '{full_text.strip()[:100]}...'")

        try:
            line_infos = self._get_line_ranges(context.original_image, bbox, context.coords)
            print(f"Detected lines: {len(line_infos)}")
            if not line_infos: raise ValueError("No lines detected.")

            raw_chars = self._detect_raw_characters(context.original_image, line_infos, bbox, context.coords)

            # Removed old color-based bullet heuristic - replaced with intelligent detection above

            # Use enhanced character alignment system for better reliability
            corrected_chars, can_align = self._enhanced_character_alignment(raw_chars, full_text, context.coords)
            context.add_characters(raw_chars, corrected_chars)

            print(f"Using char-by-char styling (enhanced alignment): {can_align}")

            if not can_align:
                # Fallback to line-based rendering
                text_runs = self._get_text_runs_by_line(all_spans, line_infos, context.coords, elem.get("type"))
                context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs})
                return

            final_styles = self._determine_character_styles(corrected_chars, context.coords, elem.get("type"))
            final_styles = self._normalize_font_sizes(final_styles)
            final_styles = self._normalize_colors(final_styles)

            # Determine if single-line based on character analysis BEFORE adding to context
            line_indices = {char.line_index for char in final_styles}
            is_single_line = len(line_indices) <= 1

            text_runs = self._get_text_runs_by_char(full_text, final_styles)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})

        except Exception:
            # Broad exception fallback
            text_runs = self._get_text_runs_from_spans(all_spans, bbox, context.original_image, context.coords, elem.get("type"))
            # Fallback check for single-line
            is_single_line = not any('\n' in run['text'] for run in text_runs)
            context.add_processed_element('text', {'bbox': bbox, 'text_runs': text_runs, 'is_single_line': is_single_line})
    
    def _update_spans_with_corrected_text(self, all_spans, original_text, corrected_text):
        """Update spans with corrected text, trying to preserve structure"""
        if len(all_spans) == 0:
            return
        
        # If text length is similar, try to map changes
        if abs(len(original_text) - len(corrected_text)) < len(original_text) * 0.3:
            # Calculate proportional distribution
            original_lengths = [len(span.get("content", "")) for span in all_spans]
            total_original_length = sum(original_lengths)
            
            if total_original_length > 0:
                corrected_pos = 0
                for i, span in enumerate(all_spans):
                    proportion = original_lengths[i] / total_original_length
                    span_length = int(len(corrected_text) * proportion)
                    
                    if i == len(all_spans) - 1:  # Last span gets remainder
                        span["content"] = corrected_text[corrected_pos:]
                    else:
                        span["content"] = corrected_text[corrected_pos:corrected_pos + span_length]
                        corrected_pos += span_length
            else:
                # Put all corrected text in the first span
                all_spans[0]["content"] = corrected_text
                for span in all_spans[1:]:
                    span["content"] = ""
    
    async def _collect_text_content(self, context, elem):
        """Collect text content for batch AI correction"""
        all_spans = [s for l in elem.get("lines", []) for s in l.get("spans", [])] if "lines" in elem else elem.get("spans", [])
        
        if not all_spans:
            text_content = elem.get("text", "")
            if text_content and text_content.strip():
                elem_type = elem.get("type", "text")
                context_info = f"Element type: {elem_type}"
                context.add_text_for_batch_correction(text_content, context_info, elem)
            return

        # Extract full text from spans
        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in all_spans])
        if not full_text.strip():
            return
        
        # Add text for batch correction
        elem_type = elem.get("type", "text")
        context_info = f"Element type: {elem_type}"
        if elem_type in ["title", "header"]:
            context_info += " (heading/title text)"
        elif elem_type in ["caption", "footnote"]:
            context_info += " (caption/footnote text)"
        
        context.add_text_for_batch_correction(full_text, context_info, elem)
    
    async def _collect_list_texts(self, context, elem):
        """Collect list text content for batch AI correction"""
        for block in elem.get("blocks", []):
            spans = [s for l in block.get("lines", []) for s in l.get("spans", [])] if "lines" in block else block.get("spans", [])
            if spans:
                full_text = "".join([s.get("content", "").replace('\\%', '%') for s in spans])
                if full_text.strip():
                    context_info = "List item text"
                    context.add_text_for_batch_correction(full_text, context_info, block)

    def _get_text_runs_by_char(self, full_text, final_styles):
        """Generates styled text runs from character-by-character analysis."""
        runs = []
        style_iter = iter(final_styles)
        last_style = None
        for char in full_text:
            # Use language-aware font selection as default
            optimal_font = self._select_optimal_font_for_language(full_text)
            font_info = {'name': optimal_font}
            
            if char not in " \n":
                style = next(style_iter, None)
                if style:
                    font_info['size'] = Pt(style.font_size)
                    font_info['color'] = RGBColor(*style.color)
                    font_info['bold'] = style.bold
                    last_style = style
            elif last_style:
                font_info['size'] = Pt(last_style.font_size)
                font_info['color'] = RGBColor(*last_style.color)
                font_info['bold'] = last_style.bold
            runs.append({'text': char, 'font': font_info})
        return runs

    def _get_text_runs_by_line(self, all_spans, line_infos, coords, elem_type):
        """Generates styled text runs using line-based analysis as a fallback with language-aware fonts."""
        runs = []
        span_idx = 0
        
        # Extract full text for language detection
        full_text = "".join([s.get("content", "") for s in all_spans])
        optimal_font = self._select_optimal_font_for_language(full_text)
        
        for i, info in enumerate(line_infos):
            line_range = info['range']
            line_spans = []
            while span_idx < len(all_spans):
                span = all_spans[span_idx]
                sbbox = span.get("bbox")
                if sbbox and sbbox[1] < line_range[1] and sbbox[3] > line_range[0]:
                    line_spans.append(span)
                    span_idx += 1
                else:
                    break
            if not line_spans: continue

            line_text = "".join([s.get("content", "").replace('\\%', '%') for s in line_spans])
            if not line_text.strip() and i < len(line_infos) - 1: line_text += "\n"

            font_size_pts = (line_range[1] - line_range[0]) * coords['scale_y']
            font_info = {
                'name': optimal_font,
                'color': RGBColor(*info['color']),
                'size': Pt(int(max(font_size_pts, 6.0))),
                'bold': elem_type == "title"
            }
            runs.append({'text': line_text, 'font': font_info})
            if i < len(line_infos) - 1 and not line_text.endswith('\n'):
                runs.append({'text': '\n', 'font': font_info}) # Keep consistent font for newline
        return runs

    def _get_text_runs_from_spans(self, spans, bbox, page_image, coords, elem_type=None):
        """Generates a single styled text run as a last-resort fallback."""
        if not spans: return []
        font_size_pts = (bbox[3] - bbox[1]) * coords['scale_y']
        full_text = "".join([s.get("content", "").replace('\\%', '%') for s in spans])
        bg_color = extract_background_color(page_image, bbox)
        color, _, _ = extract_font_color(page_image, bbox, bg_color)
        # Use language-aware font selection
        optimal_font = self._select_optimal_font_for_language(full_text)
        
        font_info = {
            'name': optimal_font,
            'size': Pt(int(font_size_pts)),
            'bold': elem_type == "title",
            'color': RGBColor(*color)
        }
        return [{'text': full_text, 'font': font_info}]

    def _render_text_from_data(self, slide, text_data, existing_elements=None):
        """Renders a text element from processed data onto a slide with collision detection."""
        bbox = text_data['bbox']
        text_runs = text_data.get('text_runs', [])
        is_single_line = text_data.get('is_single_line', False)

        # Extract text content for optimal sizing
        text_content = ''.join(run.get('text', '') for run in text_runs)

        # Calculate optimal text box size with collision detection
        if is_single_line and existing_elements is not None:
            render_bbox = self._calculate_optimal_textbox_size(
                bbox, text_content, existing_elements, is_single_line=True
            )
            # Log if we prevented an overlap
            if render_bbox != bbox:
                x1, y1, x2, y2 = bbox
                new_x1, new_y1, new_x2, new_y2 = render_bbox
                expansion = ((new_x2 - new_x1) / (x2 - x1) - 1) * 100
                print(f"   Smart text sizing: {expansion:.1f}% expansion (collision-aware)")
        elif is_single_line:
            # Fallback to old behavior if existing_elements not available
            x1, y1, x2, y2 = bbox
            width = x2 - x1
            new_x2 = x1 + width * 1.2
            render_bbox = [x1, y1, new_x2, y2]
            print(f"   Fallback text sizing: 20% expansion (no collision detection)")
        else:
            render_bbox = bbox

        txBox = self._create_textbox(slide, render_bbox, self.coords_for_render)
        tf = txBox.text_frame
        tf.clear()
        tf.margin_bottom = tf.margin_top = tf.margin_left = tf.margin_right = Pt(0)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        for run_data in text_runs:
            run = p.add_run()
            run.text = run_data['text']
            font = run.font
            font_info = run_data.get('font', {})
            
            # Use language-aware font selection if no specific font is provided
            if 'name' not in font_info or font_info['name'] == "Microsoft YaHei":
                # Extract text for language detection
                text_content = run.text
                optimal_font = self._select_optimal_font_for_language(text_content)
                font.name = optimal_font
                if optimal_font != "Microsoft YaHei":
                    print(f"   Language-aware font: {optimal_font} for text '{text_content[:20]}...'")
            else:
                font.name = font_info['name']
            
            if 'size' in font_info: font.size = font_info['size']
            if 'color' in font_info: font.color.rgb = font_info['color']
            if 'bold' in font_info: font.bold = font_info['bold']

    def _add_picture_from_bbox(self, slide, bbox, page_image, coords, text_elements):
        if not bbox: return
        x1, y1, x2, y2 = bbox;
        left, top, w, h = Pt(x1 * coords['scale_x']), Pt(y1 * coords['scale_y']), Pt((x2 - x1) * coords['scale_x']), Pt(
            (y2 - y1) * coords['scale_y'])
        px_box = [int(x1 * coords['img_w'] / coords['json_w']), int(y1 * coords['img_h'] / coords['json_h']),
                  int(x2 * coords['img_w'] / coords['json_w']), int(y2 * coords['img_h'] / coords['json_h'])]
        crop = page_image[px_box[1]:px_box[3], px_box[0]:px_box[2]].copy()

        # This cleanup logic is now less critical due to the global background inpainting,
        # but can still be useful for images that contain text not defined as a separate text element.
        for txt_e in text_elements:
            txt_box = txt_e.get("bbox")
            if txt_box and self._get_bbox_intersection(bbox, txt_box):
                px_txt_box = [int(v * (
                    coords['img_w'] / coords['json_w'] if i % 2 == 0 else coords['img_h'] / coords['json_h'])) for
                              i, v in enumerate(txt_box)]
                inter = self._get_bbox_intersection(px_box, px_txt_box)
                if inter:
                    local_inter = [inter[0] - px_box[0], inter[1] - px_box[1], inter[2] - px_box[0],
                                   inter[3] - px_box[1]]
                    fill_bbox_with_bg(crop, local_inter)

        if crop.size > 0:
            path = f"temp_crop_img_{x1}_{y1}.png";
            cv2.imwrite(path, cv2.cvtColor(crop, cv2.COLOR_RGB2BGR))
            slide.shapes.add_picture(path, left, top, w, h);
            os.remove(path)

    def _process_image(self, context, elem, text_elements):
        context.add_element_bbox_for_cleanup(elem.get("bbox"))

        if "blocks" in elem and elem["blocks"]:
            image_block_bbox = None
            for block in elem["blocks"]:
                if block.get("type") == "image_body" or (block.get("spans") and block["spans"][0].get("type") == "image"):
                    image_block_bbox = block.get("bbox")
                    break

            # Add the main image part to the render queue
            img_bbox_to_render = image_block_bbox or elem.get("bbox")
            context.add_processed_element('image', {'bbox': img_bbox_to_render, 'text_elements': text_elements})

            # Process any text blocks within the image element
            for block in elem["blocks"]:
                if block.get("type") == "image_caption":
                    self._process_text(context, block)
        else:
            # Simple image
            context.add_processed_element('image', {'bbox': elem.get("bbox"), 'text_elements': text_elements})

    def _process_list(self, context, elem):
        for block in elem.get("blocks", []):
            # Apply intelligent bullet point detection
            spans = [s for l in block.get("lines", []) for s in l.get("spans", [])] if "lines" in block else block.get("spans", [])
            if spans:
                spans.sort(key=lambda s: (s.get("bbox", [0,0,0,0])[1], s.get("bbox", [0,0,0,0])[0]))
                
                # Get full block text for intelligent analysis
                block_text = "".join([s.get("content", "") for s in spans])
                corrected_text = self._detect_bullet_points_intelligently(block_text, "list", context)
                
                if corrected_text != block_text:
                    # If text was modified, update first span with corrected text
                    spans[0]["content"] = corrected_text
                    # Clear other spans to avoid duplication
                    for span in spans[1:]:
                        span["content"] = ""
                else:
                    # Fallback to simple bullet prepending if no intelligence detected
                    spans[0]["content"] = "• " + spans[0].get("content", "").lstrip(' ·-*•')

            # Re-assign modified spans back to the block before processing
            if "lines" in block and block["lines"]:
                block["lines"][0]["spans"] = spans
            else:
                block["spans"] = spans

            self._process_text(context, block)
    
    async def _process_list_async(self, context, elem):
        """Async version of list processing with AI text correction and intelligent bullet detection"""
        for block in elem.get("blocks", []):
            # Get original text for cache lookup
            spans = [s for l in block.get("lines", []) for s in l.get("spans", [])] if "lines" in block else block.get("spans", [])
            if spans:
                # Get original text before bullet modification
                original_text = "".join([s.get("content", "").replace('\\%', '%') for s in spans])
                corrected_text = None
                
                # Check if we have corrected text in cache
                if self.enable_ai_correction and self.ai_initialized:
                    active_provider = ai_config.get_active_provider()
                    corrected_text = global_cache_manager.get_corrected_text(original_text, active_provider)
                    
                if corrected_text and corrected_text != original_text:
                    # Update spans with corrected text
                    if len(spans) == 1:
                        spans[0]["content"] = corrected_text
                    else:
                        # Distribute corrected text across spans
                        self._update_spans_with_corrected_text(spans, original_text, corrected_text)
                    working_text = corrected_text
                else:
                    working_text = original_text
                
                # Apply intelligent bullet point detection
                bullet_corrected_text = self._detect_bullet_points_intelligently(working_text, "list", context)
                
                if bullet_corrected_text != working_text:
                    # Update spans with bullet-corrected text
                    spans.sort(key=lambda s: (s.get("bbox", [0,0,0,0])[1], s.get("bbox", [0,0,0,0])[0]))
                    spans[0]["content"] = bullet_corrected_text
                    # Clear other spans to avoid duplication
                    for span in spans[1:]:
                        span["content"] = ""
                else:
                    # Fallback to simple bullet prepending if no intelligence detected
                    spans.sort(key=lambda s: (s.get("bbox", [0,0,0,0])[1], s.get("bbox", [0,0,0,0])[0]))
                    first_content = spans[0].get("content", "").lstrip(' ·-*•')
                    spans[0]["content"] = "• " + first_content

            # Re-assign modified spans back to the block before processing
            if "lines" in block and block["lines"]:
                block["lines"][0]["spans"] = spans
            else:
                block["spans"] = spans

            await self._process_text_async(context, block)

    def _process_element(self, context, elem, all_text_elements):
        cat = elem.get("type", "text")
        if cat == "list":
            self._process_list(context, elem)
        elif cat in ["text", "title", "caption", "footnote", "footer", "header", "page_number"]:
            self._process_text(context, elem)
        elif cat in ["image", "table", "formula", "figure"]:
            self._process_image(context, elem, all_text_elements)
    
    async def _collect_texts_for_batch(self, context, elem, all_text_elements):
        """Collect texts from elements for batch AI correction"""
        cat = elem.get("type", "text")
        if cat == "list":
            await self._collect_list_texts(context, elem)
        elif cat in ["text", "title", "caption", "footnote", "footer", "header", "page_number"]:
            await self._collect_text_content(context, elem)
        # Images don't need text correction

    async def _process_element_async(self, context, elem, all_text_elements):
        """Async version of element processing with AI text correction"""
        cat = elem.get("type", "text")
        if cat == "list":
            await self._process_list_async(context, elem)
        elif cat in ["text", "title", "caption", "footnote", "footer", "header", "page_number"]:
            await self._process_text_async(context, elem)
        elif cat in ["image", "table", "formula", "figure"]:
            self._process_image(context, elem, all_text_elements)  # Images don't need AI correction

    def process_page(self, slide, elements, page_image, page_size=None, page_index=0, debug_images=False):
        """Synchronous version for backward compatibility"""
        return asyncio.run(self.process_page_async(slide, elements, page_image, page_size, page_index, debug_images))
    
    async def process_page_async(self, slide, elements, page_image, page_size=None, page_index=0, debug_images=False):
        """Async version with AI text correction support"""
        self.debug_images = debug_images
        img_h, img_w = page_image.shape[:2]
        json_w, json_h = page_size if page_size and all(page_size) else (img_w * 72 / 300, img_h * 72 / 300)
        w_pts, h_pts = self.cap_size(json_w, json_h)
        self.prs.slide_width, self.prs.slide_height = Pt(w_pts), Pt(h_pts)
        coords = {'scale_x': w_pts / json_w, 'scale_y': h_pts / json_h, 'img_w': img_w, 'img_h': img_h,
                  'json_w': json_w, 'json_h': json_h}
        self.coords_for_render = coords # Store for render phase

        context = PageContext(page_image, coords, slide, page_index)

        text_types = ["list", "text", "title", "caption", "footnote", "footer", "header", "page_number"]
        all_text_elements = [e for e in elements if e.get("type", "text") in text_types]

        # Phase 1: Analyze and populate context.
        # First, register bboxes for background cleanup.
        # Erase all elements, or just discarded blocks if watermark removal is on.
        for elem in elements:
            is_discarded = elem.get('is_discarded', False)
            if not is_discarded or (is_discarded and self.remove_watermark):
                context.add_element_bbox_for_cleanup(elem.get("bbox"))

        # Second, collect texts for batch AI correction
        if self.enable_ai_correction and self.ai_initialized:
            print(f"Collecting texts for batch AI correction on page {page_index + 1}...")
            for elem in elements:
                # If watermark removal is on, skip discarded blocks.
                if elem.get('is_discarded') and self.remove_watermark:
                    continue
                await self._collect_texts_for_batch(context, elem, all_text_elements)
            
            # Apply batch text correction to all collected texts
            print(f"Applying batch AI text correction to {len(context.pending_texts)} texts...")
            batch_success = await context.apply_ai_batch_text_correction()
            if batch_success:
                print(f"Batch AI text correction completed for page {page_index + 1}")
            else:
                print(f"Batch AI text correction failed for page {page_index + 1}, will use individual correction as fallback")

        # Third, process elements with corrected texts
        for elem in elements:
            # If watermark removal is on, skip processing/rendering discarded blocks.
            if elem.get('is_discarded') and self.remove_watermark:
                continue
            await self._process_element_async(context, elem, all_text_elements)

        # Phase 2: Render from context
        context.render_to_slide(self)

        # Phase 3: Generate debug output if enabled
        if self.debug_images:
            context.generate_debug_images(page_index, self)
        
        # Return correction summary
        correction_summary = context.get_correction_summary()
        return correction_summary.get("total_corrections", 0)

    def save(self):
        self.prs.save(self.output_path)


def _get_pages_to_process(total_pages, total_images, page_selection):
    """
    Determine which pages to process based on page selection criteria.
    
    Args:
        total_pages: Total number of pages in the document
        total_images: Total number of images available
        page_selection: Page selection configuration dict
    
    Returns:
        List of 0-based page indices to process
    """
    # Use the smaller of pages or images available
    available_pages = min(total_pages, total_images)
    
    if not page_selection or page_selection.get("mode") == "all":
        return list(range(available_pages))
    
    mode = page_selection.get("mode")
    
    if mode == "single":
        page_num = page_selection.get("page", 1)
        # Convert from 1-based to 0-based index
        page_index = page_num - 1
        if 0 <= page_index < available_pages:
            return [page_index]
        else:
            print(f"Warning: Requested page {page_num} is out of range (1-{available_pages}). Processing page 1 instead.")
            return [0] if available_pages > 0 else []
    
    elif mode == "range":
        from_page = page_selection.get("from", 1)
        to_page = page_selection.get("to", 1)
        
        # Convert from 1-based to 0-based indices
        from_index = max(0, from_page - 1)
        to_index = min(available_pages - 1, to_page - 1)
        
        if from_index > to_index or from_index >= available_pages:
            print(f"Warning: Invalid page range {from_page}-{to_page} (available: 1-{available_pages}). Processing page 1 instead.")
            return [0] if available_pages > 0 else []
        
        # Adjust range if it exceeds available pages
        if from_page > available_pages:
            print(f"Warning: Start page {from_page} exceeds available pages ({available_pages}). Processing page 1 instead.")
            return [0]
        
        if to_page > available_pages:
            print(f"Warning: End page {to_page} exceeds available pages ({available_pages}). Processing pages {from_page}-{available_pages}.")
            to_index = available_pages - 1
        
        return list(range(from_index, to_index + 1))
    
    # Fallback to all pages
    return list(range(available_pages))


def convert_mineru_to_ppt(json_path, input_path, output_ppt_path, remove_watermark=True, debug_images=False, enable_ai_correction=None, page_selection=None):
    """
    Convert MinerU JSON and input file to PowerPoint presentation.
    
    Args:
        json_path: Path to MinerU JSON file
        input_path: Path to PDF or image file
        output_ppt_path: Path for output PowerPoint file
        remove_watermark: Remove watermark elements
        debug_images: Generate debug images
        enable_ai_correction: Enable AI text correction (None = use config default)
        page_selection: Dict with page selection info {"mode": "all|single|range", "page": int, "from": int, "to": int}
    """
    
    # Run the async conversion function
    return asyncio.run(_convert_mineru_to_ppt_async(
        json_path, input_path, output_ppt_path, 
        remove_watermark, debug_images, enable_ai_correction, page_selection
    ))


async def _convert_mineru_to_ppt_async(json_path, input_path, output_ppt_path, remove_watermark=True, debug_images=False, enable_ai_correction=None, page_selection=None):
    """Async implementation of the conversion process"""
    from .utils import pdf_to_images
    DPI = 300

    if debug_images:
        if os.path.exists("tmp"):
            shutil.rmtree("tmp")
        os.makedirs("tmp")

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Generate page images from either a PDF or a single image file
    if input_path.lower().endswith('.pdf'):
        images = pdf_to_images(input_path, dpi=DPI)
    else:
        try:
            img = Image.open(input_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            images = [np.array(img)]
        except Exception as e:
            raise IOError(f"Failed to load image file: {input_path} - {e}")

    gen = PPTGenerator(output_ppt_path, remove_watermark=remove_watermark, enable_ai_correction=enable_ai_correction)
    
    # Initialize AI services if enabled
    if gen.enable_ai_correction:
        print("Initializing AI services for text correction...")
        ai_success = await gen.initialize_ai_services()
        if ai_success:
            print(f"AI text correction enabled with provider: {ai_config.get_active_provider()}")
        else:
            print("AI text correction disabled - no services available")
            gen.enable_ai_correction = False
    
    pages = data if isinstance(data, list) else next(
        (data[k] for k in ["pdf_info", "pages"] if k in data and isinstance(data[k], list)), [data])
    print(f"[CLEANUP] Found {len(pages)} pages.")
    
    # Determine which pages to process based on page_selection
    pages_to_process = _get_pages_to_process(len(pages), len(images), page_selection)
    print(f"[PAGE SELECTION] Processing {len(pages_to_process)} page(s): {pages_to_process}")
    
    total_corrections = 0
    for slide_index, page_index in enumerate(pages_to_process):
        page_data = pages[page_index]
        page_img = images[page_index].copy()
        print(f"Processing page {page_index + 1}/{len(pages)} (slide {slide_index + 1}/{len(pages_to_process)})...")
        
        if slide_index == 0: gen.set_slide_size(page_img.shape[1], page_img.shape[0], dpi=DPI)
        slide = gen.add_slide()

        elements = []
        for key in ["para_blocks", "images", "tables"]:
            for item in page_data.get(key, []):
                item['is_discarded'] = False
                elements.append(item)
        for item in page_data.get("discarded_blocks", []):
            item['is_discarded'] = True
            elements.append(item)

        page_size = page_data.get("page_size") or (page_data.get("page_info", {}).get("width"),
                                                   page_data.get("page_info", {}).get("height"))
        
        # Process page with AI correction
        page_corrections = await gen.process_page_async(slide, elements, page_img, page_size=page_size, page_index=page_index, debug_images=debug_images)
        total_corrections += page_corrections
    
    gen.save()
    
    # Print AI correction summary
    if gen.enable_ai_correction and total_corrections > 0:
        print(f"AI Text Correction Summary: {total_corrections} corrections applied")
        active_provider = ai_config.get_active_provider()
        if active_provider:
            print(f"Provider used: {active_provider}")
    
    print(f"Saved to {output_ppt_path}")
